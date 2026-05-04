from pathlib import Path
from uuid import uuid4

from fastapi import UploadFile
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.config import UPLOAD_DIR


def save_uploaded_file(file: UploadFile) -> Path:
    original_name = Path(file.filename or "upload").name
    saved_name = f"{uuid4().hex}_{original_name}"
    file_path = UPLOAD_DIR / saved_name

    with file_path.open("wb") as buffer:
        buffer.write(file.file.read())

    return file_path


def extract_text_from_file(file_path: str) -> dict:
    path = Path(file_path)
    file_type = path.suffix.lower().lstrip(".")

    parsers = {
        "pptx": extract_text_from_pptx,
        "docx": extract_text_from_docx,
        "pdf": extract_text_from_pdf,
    }

    if file_type not in parsers:
        return {
            "file_type": "unknown",
            "text": "",
            "characters": 0,
            "status": "unsupported",
            "error": "Поддерживаются только PPTX, DOCX и PDF",
        }

    try:
        text = parsers[file_type](str(path))
        return {
            "file_type": file_type,
            "text": text,
            "characters": len(text),
            "status": "ok",
        }
    except Exception as error:
        return {
            "file_type": file_type,
            "text": "",
            "characters": 0,
            "status": "error",
            "error": str(error),
        }


def extract_text_from_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    parts = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"--- СЛАЙД {slide_number} ---")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())

    return "\n\n".join(part for part in parts if part)


def extract_text_from_docx(file_path: str) -> str:
    document = Document(file_path)
    parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

    return "\n\n".join(parts)


def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    parts = []

    for page_number, page in enumerate(reader.pages, start=1):
        parts.append(f"--- СТРАНИЦА {page_number} ---")
        page_text = page.extract_text() or ""
        if page_text.strip():
            parts.append(page_text.strip())

    return "\n\n".join(parts)
